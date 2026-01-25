#!/usr/bin/env python3
"""Gera uma apresentacao (PPTX) em estilo inspirado no Claude.

Entrada:
- dados/pesquisas_eleitorais.csv
- dados/eventos_timeline.json

Saida:
- outputs/apresentacao_claude/graficos/*.png
- outputs/apresentacao_claude/Celina_Leao_2024-2025_Analise_ClaudeStyle.pptx
- outputs/apresentacao_claude/deck_marp.md
"""

from __future__ import annotations

import json
import math
import os
import sys
import base64
from dataclasses import dataclass
from pathlib import Path
from typing import Any, cast


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import numpy as np
import pandas as pd
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
from PIL import Image


from presentation_sdk.python.tema_anthropic import carregar_tokens, aplicar_tema_matplotlib


OUT_DIR = ROOT / "outputs" / "apresentacao_claude"
CHARTS_DIR = OUT_DIR / "graficos"


@dataclass(frozen=True)
class Metricas:
    total_pesquisas: int
    total_institutos: int
    media_geral: float
    media_pre: float
    media_pos: float
    lift_pp: float
    lift_rel: float
    efeito_d: float
    minimo: float
    maximo: float
    data_min: str
    data_max: str


def carregar_dados() -> tuple[Any, dict[str, Any]]:
    df = pd.read_csv(ROOT / "dados" / "pesquisas_eleitorais.csv")
    df["data_pesquisa"] = pd.to_datetime(df["data_pesquisa"])
    df["data_divulgacao"] = pd.to_datetime(df["data_divulgacao"])

    eventos = json.loads((ROOT / "dados" / "eventos_timeline.json").read_text(encoding="utf-8"))
    return df, eventos


def calcular_metricas(df: Any) -> Metricas:
    marco_mkt = pd.Timestamp("2025-01-01")
    pre = df[df["data_pesquisa"] < marco_mkt]
    pos = df[df["data_pesquisa"] >= marco_mkt]

    media_pre = float(pre["celina_leao"].astype(float).mean())
    media_pos = float(pos["celina_leao"].astype(float).mean())
    lift = media_pos - media_pre
    lift_rel = (lift / media_pre) if media_pre else 0.0

    # Efeito (Cohen's d) com desvio padrao combinado
    sd_pre = float(pre["celina_leao"].astype(float).std(ddof=1)) if len(pre) > 1 else 0.0
    sd_pos = float(pos["celina_leao"].astype(float).std(ddof=1)) if len(pos) > 1 else 0.0
    pooled = 0.0
    if len(pre) > 1 and len(pos) > 1:
        pooled = math.sqrt(((len(pre) - 1) * sd_pre**2 + (len(pos) - 1) * sd_pos**2) / (len(pre) + len(pos) - 2))
    d = (lift / pooled) if pooled else 0.0

    cel = df["celina_leao"].astype(float)
    idx_min = int(cel.idxmin())
    idx_max = int(cel.idxmax())
    data_min = pd.Timestamp(df.loc[idx_min, "data_divulgacao"]).strftime("%b/%Y")
    data_max = pd.Timestamp(df.loc[idx_max, "data_divulgacao"]).strftime("%b/%Y")

    return Metricas(
        total_pesquisas=int(len(df)),
        total_institutos=int(df["instituto"].nunique()),
        media_geral=float(cel.mean()),
        media_pre=media_pre,
        media_pos=media_pos,
        lift_pp=lift,
        lift_rel=lift_rel,
        efeito_d=d,
        minimo=float(cel.loc[idx_min]),
        maximo=float(cel.loc[idx_max]),
        data_min=data_min,
        data_max=data_max,
    )


def _save(fig: Figure, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=220, bbox_inches="tight")
    plt.close(fig)


def grafico_evolucao(df: Any, tokens: Any) -> Path:
    df1 = df[df["cenario"] == 1].sort_values(by="data_pesquisa").copy()
    marco = pd.Timestamp("2025-01-01")

    fig, ax = plt.subplots(figsize=(13, 6.4))
    aplicar_tema_matplotlib(tokens)

    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59"]) or ["#8fa6d3"]
    c1 = accents[0]
    c2 = accents[1] if len(accents) > 1 else accents[0]

    ax.axvspan(marco, df1["data_pesquisa"].max(), color=c1, alpha=0.08)  # type: ignore[arg-type]
    ax.plot(df1["data_pesquisa"], df1["celina_leao"], marker="o", linewidth=2.8, color=c1)
    ax.plot(df1["data_pesquisa"], df1["segundo_colocado"], marker="o", linewidth=2.2, color=c2, alpha=0.9)

    ax.axvline(marco, color=tokens.colors.get("muted", "#6c5b68"), linestyle="--", linewidth=1.5, alpha=0.8)  # type: ignore[arg-type]
    ax.text(marco, ax.get_ylim()[1] * 0.98, "  Jan/2025\n  inicio MKT", va="top", fontsize=10)  # type: ignore[arg-type]

    ax.set_title("Evolucao da intencao de voto (cenario 1)")
    ax.set_xlabel("Data")
    ax.set_ylabel("%")
    ax.legend(["Celina Leao", "2o colocado"], loc="upper left")

    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b/%y"))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    fig.autofmt_xdate(rotation=0)

    out = CHARTS_DIR / "01_evolucao_claude.png"
    _save(fig, out)
    return out


def grafico_pre_pos(df: Any, m: Metricas, tokens: Any) -> Path:
    marco = pd.Timestamp("2025-01-01")
    pre = df[df["data_pesquisa"] < marco]
    pos = df[df["data_pesquisa"] >= marco]
    sd_pre = float(pre["celina_leao"].std(ddof=1)) if len(pre) > 1 else 0.0
    sd_pos = float(pos["celina_leao"].std(ddof=1)) if len(pos) > 1 else 0.0

    fig, ax = plt.subplots(figsize=(10.5, 6.0))
    aplicar_tema_matplotlib(tokens)
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76"]) or ["#8fa6d3"]

    x = np.arange(2)
    vals = [m.media_pre, m.media_pos]
    errs = [sd_pre, sd_pos]
    bars = ax.bar(x, vals, yerr=errs, capsize=8, color=[tokens.colors.get("muted", "#6c5b68"), accents[0]])

    for i, b in enumerate(bars):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.9, f"{vals[i]:.1f}%", ha="center", fontsize=12)

    ax.set_xticks(x, ["Pre-MKT\n(jul-dez/2024)", "Pos-MKT\n(2025)"])
    ax.set_ylabel("Media (%)")
    ax.set_title(f"Impacto do marketing: +{m.lift_pp:.1f} p.p. (≈ {m.lift_rel*100:.0f}% )")
    ax.set_ylim(0, max(vals) + 15)

    out = CHARTS_DIR / "02_pre_pos_claude.png"
    _save(fig, out)
    return out


def grafico_tendencia(df: Any, tokens: Any) -> Path:
    df1 = df[df["cenario"] == 1].sort_values(by="data_pesquisa").copy()
    marco = pd.Timestamp("2025-01-01")
    df1["dias"] = (df1["data_pesquisa"] - df1["data_pesquisa"].min()).dt.days

    pre = df1[df1["data_pesquisa"] < marco]
    pos = df1[df1["data_pesquisa"] >= marco]

    fig, ax = plt.subplots(figsize=(11.5, 6.2))
    aplicar_tema_matplotlib(tokens)
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76"]) or ["#8fa6d3"]

    ax.scatter(pre["dias"], pre["celina_leao"], s=90, color=tokens.colors.get("muted", "#6c5b68"), alpha=0.9)
    ax.scatter(pos["dias"], pos["celina_leao"], s=90, color=accents[0], alpha=0.95)

    def fit_line(dfi, color, label):
        if len(dfi) < 2:
            return
        z = np.polyfit(dfi["dias"], dfi["celina_leao"], 1)
        p = np.poly1d(z)
        xs = np.linspace(dfi["dias"].min(), dfi["dias"].max(), 100)
        ax.plot(xs, p(xs), color=color, linewidth=2.6, label=f"{label} (slope {z[0]:.3f})")

    fit_line(pre, tokens.colors.get("muted", "#6c5b68"), "Pre")
    fit_line(pos, accents[0], "Pos")

    dias_marco = int((marco - df1["data_pesquisa"].min()).days)
    ax.axvline(dias_marco, color=tokens.colors.get("stroke", "#d8d2c9"), linestyle="--", linewidth=1.5)
    ax.text(dias_marco + 8, ax.get_ylim()[1] * 0.98, "Jan/2025", va="top", fontsize=10)

    ax.set_title("Mudanca de tendencia (cenario 1)")
    ax.set_xlabel("Dias desde a 1a pesquisa")
    ax.set_ylabel("%")
    ax.legend(loc="upper left")

    out = CHARTS_DIR / "03_tendencia_quebra_claude.png"
    _save(fig, out)
    return out


def grafico_institutos(df: Any, tokens: Any) -> Path:
    df_ultimo = df.sort_values("data_pesquisa").groupby("instituto").last().reset_index()
    df_ultimo = df_ultimo.sort_values("celina_leao")

    fig, ax = plt.subplots(figsize=(11.5, 6.2))
    aplicar_tema_matplotlib(tokens)
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59", "#72797c"]) or ["#8fa6d3"]

    colors = [accents[i % len(accents)] for i in range(len(df_ultimo))]
    bars = ax.barh(df_ultimo["instituto"], df_ultimo["celina_leao"], color=colors, height=0.55)
    for b in bars:
        ax.text(b.get_width() + 0.8, b.get_y() + b.get_height() / 2, f"{b.get_width():.1f}%", va="center", fontsize=11)

    ax.set_title("Ultima medicao por instituto")
    ax.set_xlabel("%")

    out = CHARTS_DIR / "04_institutos_claude.png"
    _save(fig, out)
    return out


def grafico_competicao(df: Any, tokens: Any) -> Path:
    ultimo = df.sort_values("data_pesquisa").iloc[-1]
    nomes = ["Celina Leao", str(ultimo["nome_segundo"]), str(ultimo["nome_terceiro"])]
    vals = [float(ultimo["celina_leao"]), float(ultimo["segundo_colocado"]), float(ultimo["terceiro_colocado"])]
    pares = [(n, v) for n, v in zip(nomes, vals) if n and n != "N/A" and v > 0]
    pares.sort(key=lambda x: x[1])

    fig, ax = plt.subplots(figsize=(11.5, 6.2))
    aplicar_tema_matplotlib(tokens)
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59"]) or ["#8fa6d3"]

    colors = [accents[i % len(accents)] for i in range(len(pares))]
    bars = ax.barh([p[0] for p in pares], [p[1] for p in pares], color=colors, height=0.6)
    for b in bars:
        ax.text(b.get_width() + 0.8, b.get_y() + b.get_height() / 2, f"{b.get_width():.1f}%", va="center", fontsize=11)

    ax.set_title(f"Cenario competitivo (ultima divulgacao: {ultimo['data_divulgacao'].strftime('%d/%m/%Y')})")
    ax.set_xlabel("%")
    out = CHARTS_DIR / "05_competicao_claude.png"
    _save(fig, out)
    return out


def grafico_timeline(df: Any, eventos: dict[str, Any], tokens: Any) -> Path:
    df1 = df[df["cenario"] == 1].sort_values(by="data_pesquisa").copy()
    marco = pd.Timestamp("2025-01-01")

    fig, ax = plt.subplots(figsize=(13.5, 6.5))
    aplicar_tema_matplotlib(tokens)
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59"]) or ["#8fa6d3"]

    ax.plot(df1["data_pesquisa"], df1["celina_leao"], marker="o", linewidth=2.8, color=accents[0])
    ax.axvspan(marco, df1["data_pesquisa"].max(), color=accents[0], alpha=0.08)  # type: ignore[arg-type]

    # Seleciona 8 eventos principais (prioriza muito_positivo/negativo)
    evs = []
    for e in eventos.get("eventos", []):
        d = pd.Timestamp(e["data"])
        if d < df1["data_pesquisa"].min() or d > df1["data_pesquisa"].max():
            continue
        evs.append(e)

    def score(e):
        imp = e.get("impacto")
        if imp == "muito_positivo":
            return 3
        if imp == "positivo":
            return 2
        if imp == "negativo":
            return 2
        return 1

    evs = sorted(evs, key=score, reverse=True)[:8]
    evs = sorted(evs, key=lambda e: e["data"])

    y_min, y_max = ax.get_ylim()
    y_top = y_max - (y_max - y_min) * 0.08
    direction = 1
    step = (y_max - y_min) * 0.10

    for i, e in enumerate(evs):
        d = pd.Timestamp(e["data"])
        # y do ponto mais proximo
        idx = (df1["data_pesquisa"] - d).abs().idxmin()
        y = float(df1.loc[idx, "celina_leao"])

        ax.axvline(d, color=tokens.colors.get("stroke", "#d8d2c9"), linestyle=":", linewidth=1.2)  # type: ignore[arg-type]
        yy = y_top - i * step * 0.22
        if direction < 0:
            yy = y_min + (i + 1) * step * 0.22
        direction *= -1

        ax.annotate(  # type: ignore[arg-type]
            e["titulo"],
            xy=cast(Any, (d, y)),
            xytext=cast(Any, (d, yy)),
            textcoords="data",
            fontsize=9,
            ha="left",
            va="center",
            arrowprops=dict(arrowstyle="-", color=tokens.colors.get("muted", "#6c5b68"), alpha=0.7),
            bbox=dict(boxstyle="round,pad=0.25", facecolor=tokens.colors.get("surface", "#ffffff"), edgecolor=tokens.colors.get("stroke", "#d8d2c9"), alpha=0.95),
        )

    ax.set_title("Marcos e eventos (amostra de maior impacto)")
    ax.set_xlabel("Data")
    ax.set_ylabel("%")
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b/%y"))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    fig.autofmt_xdate(rotation=0)

    out = CHARTS_DIR / "06_timeline_claude.png"
    _save(fig, out)
    return out


def grafico_waterfall(df: Any, tokens: Any) -> Path:
    df1 = df[df["cenario"] == 1].sort_values(by="data_pesquisa").copy()
    if len(df1) < 2:
        out = CHARTS_DIR / "07_waterfall_claude.png"
        fig, _ax = plt.subplots(figsize=(10, 5))
        aplicar_tema_matplotlib(tokens)
        _save(fig, out)
        return out

    labels = df1["data_divulgacao"].dt.strftime("%b/%y").tolist()
    vals = df1["celina_leao"].astype(float).tolist()
    deltas = [vals[0]] + [vals[i] - vals[i - 1] for i in range(1, len(vals))]

    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59"]) or ["#8fa6d3"]
    pos_color = accents[0]
    neg_color = accents[1] if len(accents) > 1 else tokens.colors.get("muted", "#6c5b68")
    base_color = tokens.colors.get("muted", "#6c5b68")

    fig, ax = plt.subplots(figsize=(13.0, 6.2))
    aplicar_tema_matplotlib(tokens)

    starts = [0.0]
    for i in range(1, len(deltas)):
        starts.append(starts[-1] + deltas[i - 1])

    colors = [base_color] + [pos_color if d >= 0 else neg_color for d in deltas[1:]]
    bars = ax.bar(range(len(deltas)), deltas, bottom=starts, color=colors, width=0.75)

    running = 0.0
    for i, (b, d) in enumerate(zip(bars, deltas)):
        running = starts[i] + d
        label = f"{running:.1f}%" if i == 0 else f"{d:+.1f}"
        ax.text(b.get_x() + b.get_width() / 2, running + 0.8, label, ha="center", fontsize=10)

    ax.set_xticks(range(len(labels)), labels)
    ax.set_title("Waterfall: variacao entre pesquisas (cenario 1)")
    ax.set_xlabel("Divulgacao")
    ax.set_ylabel("p.p. (acumulado)")
    for tick in ax.get_xticklabels():
        tick.set_rotation(0)

    out = CHARTS_DIR / "07_waterfall_claude.png"
    _save(fig, out)
    return out


def grafico_boxplot_institutos(df: Any, tokens: Any) -> Path:
    fig, ax = plt.subplots(figsize=(12.0, 6.2))
    aplicar_tema_matplotlib(tokens)

    insts = sorted(df["instituto"].unique().tolist())
    data = [df[df["instituto"] == inst]["celina_leao"].astype(float).values for inst in insts]
    bp = ax.boxplot(data, tick_labels=insts, patch_artist=True, widths=0.55)

    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59", "#72797c"]) or ["#8fa6d3"]
    for i, box in enumerate(bp["boxes"]):
        box.set_facecolor(accents[i % len(accents)])
        box.set_alpha(0.28)
        box.set_edgecolor(accents[i % len(accents)])
        box.set_linewidth(1.6)

    for med in bp["medians"]:
        med.set_color(tokens.colors.get("text", "#2a2d26"))
        med.set_linewidth(2.0)

    ax.set_title("Distribuicao por instituto (boxplot)")
    ax.set_xlabel("Instituto")
    ax.set_ylabel("% (Celina Leao)")
    for tick in ax.get_xticklabels():
        tick.set_rotation(0)

    out = CHARTS_DIR / "08_boxplot_institutos_claude.png"
    _save(fig, out)
    return out


def gerar_graficos(df: Any, eventos: dict[str, Any], m: Metricas, tokens: Any) -> dict[str, Path]:
    return {
        "evolucao": grafico_evolucao(df, tokens),
        "pre_pos": grafico_pre_pos(df, m, tokens),
        "tendencia": grafico_tendencia(df, tokens),
        "institutos": grafico_institutos(df, tokens),
        "competicao": grafico_competicao(df, tokens),
        "timeline": grafico_timeline(df, eventos, tokens),
        "waterfall": grafico_waterfall(df, tokens),
        "boxplot": grafico_boxplot_institutos(df, tokens),
    }


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        return (0, 0, 0)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pptx_color(hex_color: str):
    from pptx.dml.color import RGBColor

    r, g, b = _hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


def _add_bg(slide, hex_color: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    r, g, b = _hex_to_rgb(hex_color)
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_title(slide, text: str, *, left, top, width, height, color: str, size: int, bold: bool = True):
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _pptx_color(color)
    return tb


def _add_paragraphs(slide, lines: list[str], *, left, top, width, height, color: str, size: int):
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = _pptx_color(color)
    return tb


def _add_chip(slide, text: str, *, left, top, width, height, bg: str, fg: str):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _pptx_color(bg)
    shape.line.color.rgb = _pptx_color(bg)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = _pptx_color(fg)
    return shape


def _add_image_full(slide, img_path: Path, *, left, top, width, height):
    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)


def gerar_pptx(df: Any, eventos: dict[str, Any], m: Metricas, charts: dict[str, Path], tokens: Any) -> Path:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg = tokens.colors.get("bg", "#f3f0ea")
    surface = tokens.colors.get("surface", "#ffffff")
    text = tokens.colors.get("text", "#2a2d26")
    muted = tokens.colors.get("muted", "#6c5b68")
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59"]) or ["#8fa6d3"]
    accent = accents[0]

    blank = prs.slide_layouts[6]

    def slide_title(title: str, subtitle: str | None = None):
        s = prs.slides.add_slide(blank)
        _add_bg(s, bg)
        _add_chip(s, "INTEIA | Pesquisa Eleitoral DF", left=Inches(0.8), top=Inches(0.55), width=Inches(3.6), height=Inches(0.42), bg=surface, fg=muted)
        _add_title(s, title, left=Inches(0.8), top=Inches(1.25), width=Inches(12), height=Inches(1.2), color=text, size=44)
        if subtitle:
            _add_paragraphs(s, [subtitle], left=Inches(0.82), top=Inches(2.35), width=Inches(10.8), height=Inches(0.7), color=muted, size=18)
        # rodape
        tb = s.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(12), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = "Periodo analisado: Jul/2024 a Dez/2025 | Atualizado: Jan/2026"
        p.font.size = Pt(11)
        p.font.color.rgb = _pptx_color(muted)
        return s

    def slide_section(title: str, bullets: list[str]):
        s = prs.slides.add_slide(blank)
        _add_bg(s, bg)
        _add_title(s, title, left=Inches(0.8), top=Inches(0.85), width=Inches(12), height=Inches(0.9), color=text, size=34)
        _add_paragraphs(s, bullets, left=Inches(0.95), top=Inches(1.75), width=Inches(11.8), height=Inches(5.2), color=muted, size=20)
        return s

    def slide_kpis():
        s = prs.slides.add_slide(blank)
        _add_bg(s, bg)
        _add_title(s, "Resultados-chave", left=Inches(0.8), top=Inches(0.85), width=Inches(12), height=Inches(0.8), color=text, size=34)

        from pptx.enum.shapes import MSO_SHAPE

        def card(x, y, w, h, big, small):
            shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _pptx_color(surface)
            shape.line.color.rgb = _pptx_color(tokens.colors.get("stroke", "#d8d2c9"))
            tf = shape.text_frame
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = big
            p1.font.size = Pt(30)
            p1.font.bold = True
            p1.font.color.rgb = _pptx_color(text)
            p2 = tf.add_paragraph()
            p2.text = small
            p2.font.size = Pt(12)
            p2.font.color.rgb = _pptx_color(muted)

        card(Inches(0.8), Inches(1.75), Inches(3.9), Inches(1.55), f"{m.media_geral:.1f}%", "Media geral (todas as pesquisas)")
        card(Inches(4.85), Inches(1.75), Inches(3.9), Inches(1.55), f"+{m.lift_pp:.1f} p.p.", "Lift medio apos inicio do MKT")
        card(Inches(8.9), Inches(1.75), Inches(3.9), Inches(1.55), f"{m.maximo:.0f}%", f"Pico observado ({m.data_max})")
        card(Inches(0.8), Inches(3.55), Inches(3.9), Inches(1.55), f"{m.minimo:.0f}%", f"Minimo observado ({m.data_min})")
        card(Inches(4.85), Inches(3.55), Inches(3.9), Inches(1.55), f"{m.total_pesquisas}", "Registros no dataset")
        card(Inches(8.9), Inches(3.55), Inches(3.9), Inches(1.55), f"{m.total_institutos}", "Institutos (Parana, RTBD, Opiniao, Colectta)")

        # Nota interpretativa
        note = s.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(12.0), Inches(1.3))
        p = note.text_frame.paragraphs[0]
        p.text = f"Efeito estimado (Cohen's d): {m.efeito_d:.2f} | A leitura e correlacional: marketing + entregas + eventos politicos." 
        p.font.size = Pt(13)
        p.font.color.rgb = _pptx_color(muted)
        return s

    def slide_chart(title: str, img: Path, caption: str):
        s = prs.slides.add_slide(blank)
        _add_bg(s, bg)
        _add_title(s, title, left=Inches(0.8), top=Inches(0.7), width=Inches(12), height=Inches(0.7), color=text, size=28)
        _add_image_full(s, img, left=Inches(0.8), top=Inches(1.35), width=Inches(11.8), height=Inches(5.3))
        tb = s.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(12), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = _pptx_color(muted)
        return s

    # --- Conteudo ---
    slide_title(
        "Celina Leao 2024-2025",
        "Evolucao de campanha, impacto de publicidade e consolidacao de marca (DF 2026)",
    )

    slide_section(
        "Pergunta e hipotese",
        [
            "Pergunta: houve mudanca estrutural na intencao de voto apos o inicio de investimento em comunicacao (Jan/2025)?",
            "Hipotese operacional: aumento de presenca + publicidade + agenda de entregas eleva conhecimento e preferencia.",
            "Cuidado: pesquisas sao de institutos/cenarios diferentes; a leitura e de sinal agregado (nao causal estrita).",
        ],
    )

    slide_section(
        "Base e metodo",
        [
            f"Fonte: {m.total_pesquisas} registros (Jul/2024 a Dez/2025) em 4 institutos.",
            "Cortes: pre-MKT (< 01/01/2025) vs pos-MKT (>= 01/01/2025).",
            "Eixo principal de tendencia: cenario 1 (quando disponivel) para serie temporal.",
            "Medidas: medias, dispersao, efeito (Cohen's d), comparativos por instituto e cenario.",
        ],
    )

    slide_kpis()

    slide_chart(
        "Evolucao da intencao de voto",
        charts["evolucao"],
        "Serie temporal (cenario 1). Faixa sombreada: periodo apos inicio do MKT (Jan/2025).",
    )

    slide_chart(
        "Impacto: antes vs depois",
        charts["pre_pos"],
        "Comparacao de medias com barras de dispersao (desvio-padrao).",
    )

    slide_chart(
        "Mudanca de tendencia",
        charts["tendencia"],
        "Regressoes separadas sugerem inclinacao maior apos Jan/2025 (sinal de aceleracao).",
    )

    slide_chart(
        "Consistencia por instituto",
        charts["institutos"],
        "Ultima medicao de cada instituto (nao sao datas identicas; use como triangulacao).",
    )

    slide_chart(
        "Cenario competitivo",
        charts["competicao"],
        "Snapshot da ultima divulgacao do dataset (cenarios variam conforme instituto).",
    )

    slide_chart(
        "Marcos e eventos",
        charts["timeline"],
        "Eventos de alto impacto (positivo/negativo) alinhados a pontos de pesquisa.",
    )

    slide_chart(
        "Volatilidade (waterfall)",
        charts["waterfall"],
        "Decomposicao da variacao entre pesquisas (cenario 1). Ajuda a ver saltos e recuos.",
    )

    slide_chart(
        "Dispersao por instituto",
        charts["boxplot"],
        "Boxplot com distribuicao de valores por instituto (tamanho de amostra varia).",
    )

    slide_section(
        "Leituras sobre publicidade e marca",
        [
            "Marketing atua como multiplicador: aumenta lembranca, melhora consistencia da narrativa e estabiliza preferencia.",
            "Apos Jan/2025, a media sobe e o teto observado aumenta (pico em Set/2025).",
            "Marcos politicos/juridicos (ex.: absolvicacao, FCDF) funcionam como gatilhos de validacao.",
            "Risco: ruido de cenarios e entrada/saida de nomes (ex.: Fred, Arruda) desloca percentuais.",
        ],
    )

    slide_section(
        "Conclusoes",
        [
            f"Hipotese confirmada no agregado: lift medio de +{m.lift_pp:.1f} p.p. apos Jan/2025.",
            "Nao e prova causal, mas o conjunto (agenda + comunicacao + eventos) e compativel com crescimento de marca.",
            "A estrategia vencedora combina: entregas comunicaveis + ocupacao de agenda + gestao de risco juridico.",
        ],
    )

    slide_section(
        "Recomendacoes (2026)",
        [
            "Manter investimento em comunicacao e cadencia de conteudo (press + digital) com mensagens simples e repetiveis.",
            "Tratar variacao por cenario: produzir tracking proprio (mesma pergunta / mesma amostra ao longo do tempo).",
            "Criar metricas de marca (share of voice, sentimento, busca) para ligar midia -> preferencia com mais precisao.",
            "Planejar resposta rapida a riscos juridicos e ataques (janela critica de definicao de candidatura).",
        ],
    )

    out = OUT_DIR / "Celina_Leao_2024-2025_Analise_ClaudeStyle.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def gerar_marp(m: Metricas, charts: dict[str, Path]) -> Path:
    # Deck alternativo em Marp (tema anthropic_light). O usuario pode exportar com marp-cli.
    rel = lambda p: os.path.relpath(p, OUT_DIR).replace("\\", "/")
    md = f"""---
marp: true
theme: anthropic_light
paginate: true
---

# Celina Leao 2024-2025

Evolucao de campanha, impacto de publicidade e consolidacao de marca (DF 2026)

---

## Resultados-chave

- Media geral: **{m.media_geral:.1f}%**
- Lift pos-MKT: **+{m.lift_pp:.1f} p.p.** (≈ {m.lift_rel*100:.0f}%)
- Pico: **{m.maximo:.0f}%** ({m.data_max})

---

## Evolucao

![]({rel(charts['evolucao'])})

---

## Antes vs depois

![]({rel(charts['pre_pos'])})

---

## Mudanca de tendencia

![]({rel(charts['tendencia'])})

---

## Consistencia por instituto

![]({rel(charts['institutos'])})

---

## Cenario competitivo

![]({rel(charts['competicao'])})

---

## Volatilidade (waterfall)

![]({rel(charts['waterfall'])})

---

## Distribuicao por instituto

![]({rel(charts['boxplot'])})

---

## Marcos e eventos

![]({rel(charts['timeline'])})

---

## Conclusao

- Hipotese confirmada no agregado: aumento consistente apos Jan/2025.
- Leituras: comunicacao + agenda + eventos funcionam como alavancas de marca.
- Proximo passo: tracking padronizado + metricas de midia (SoV/sentimento/busca).
"""

    out = OUT_DIR / "deck_marp.md"
    out.write_text(md, encoding="utf-8")
    return out


def _b64_png(path: Path) -> str:
    data = path.read_bytes()
    return base64.b64encode(data).decode("ascii")


def gerar_html(
    df: Any,
    eventos: dict[str, Any],
    m: Metricas,
    charts: dict[str, Path],
    tokens: Any,
    *,
    download_links: bool,
    out_path: Path | None = None,
) -> Path:
    # HTML standalone com imagens embutidas (base64)
    bg = tokens.colors.get("bg", "#f3f0ea")
    surface = tokens.colors.get("surface", "#ffffff")
    text = tokens.colors.get("text", "#2a2d26")
    muted = tokens.colors.get("muted", "#6c5b68")
    stroke = tokens.colors.get("stroke", "#d8d2c9")
    accents = tokens.colors.get("accents", ["#8fa6d3", "#b67a76", "#916f59", "#72797c"]) or ["#8fa6d3"]

    # Eventos (resumo curto)
    marco_mkt = pd.Timestamp("2025-01-01")
    evs = []
    for e in eventos.get("eventos", []):
        d = pd.Timestamp(e["data"])
        if d.year < 2024 or d.year > 2025:
            continue
        evs.append((d, e.get("titulo", ""), e.get("impacto", ""), e.get("descricao", "")))
    evs.sort(key=lambda x: x[0])

    def fmt_date(d: pd.Timestamp) -> str:
        return d.strftime("%b/%Y")

    def badge(impacto: str) -> str:
        if impacto == "muito_positivo":
            return "alto+"
        if impacto == "positivo":
            return "+"
        if impacto == "negativo":
            return "-"
        return "="

    eventos_html = "".join(
        f"""<div class='event'>
  <div class='event__when'>{fmt_date(d)}</div>
  <div class='event__title'>{titulo}</div>
  <div class='event__badge'>{badge(imp)}</div>
  <div class='event__desc'>{desc}</div>
</div>"""
        for d, titulo, imp, desc in evs
    )

    def esc(s: Any) -> str:
        s = "" if s is None else str(s)
        return (
            s.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
        )

    df_tab = df.copy()
    df_tab = df_tab.sort_values(by=["data_divulgacao", "instituto", "cenario"])  # type: ignore[call-arg]
    linhas = []
    for _i, r in df_tab.iterrows():
        linhas.append(
            "".join(
                [
                    "<tr>",
                    f"<td>{esc(pd.Timestamp(r['data_divulgacao']).strftime('%d/%m/%Y'))}</td>",
                    f"<td>{esc(r['instituto'])}</td>",
                    f"<td>{esc(r['cenario'])}</td>",
                    f"<td><b>{float(r['celina_leao']):.1f}%</b></td>",
                    f"<td>{esc(r['nome_segundo'])} ({float(r['segundo_colocado']):.1f}%)</td>",
                    f"<td>{esc(r['nome_terceiro'])} ({float(r['terceiro_colocado']):.1f}%)</td>",
                    f"<td class='muted'>{esc(r.get('observacao', ''))}</td>",
                    "</tr>",
                ]
            )
        )
    tabela_html = "\n".join(linhas)

    # Charts embutidos
    chart_imgs = {k: f"data:image/png;base64,{_b64_png(v)}" for k, v in charts.items()}

    actions = """
        <a class=\"btn\" href=\"Celina_Leao_2024-2025_Analise_ClaudeStyle.pptx\" download>Baixar PPTX</a>
        <a class=\"btn\" href=\"deck_marp.md\" download>Baixar Markdown</a>
    """ if download_links else ""

    html = f"""<!doctype html>
<html lang="pt-br">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Celina Leao 2024-2025 | Analise (Claude Style)</title>
  <style>
    :root {{
      --bg: {bg};
      --surface: {surface};
      --text: {text};
      --muted: {muted};
      --stroke: {stroke};
      --a1: {accents[0]};
      --a2: {accents[1] if len(accents) > 1 else accents[0]};
      --a3: {accents[2] if len(accents) > 2 else accents[0]};
      --radius: 18px;
      --shadow: 0 18px 60px rgba(0,0,0,0.08);
    }}

    * {{ box-sizing: border-box; }}
    html, body {{ height: 100%; }}
    body {{
      margin: 0;
      color: var(--text);
      background:
        radial-gradient(1200px 600px at 15% 10%, color-mix(in oklab, var(--a1) 22%, transparent), transparent 60%),
        radial-gradient(900px 500px at 90% 15%, color-mix(in oklab, var(--a2) 18%, transparent), transparent 55%),
        linear-gradient(180deg, #faf8f4 0%, #f1ece4 100%);
      font-family: Inter, system-ui, -apple-system, "Segoe UI", Roboto, Arial, sans-serif;
      line-height: 1.4;
    }}

    .wrap {{ max-width: 1160px; margin: 0 auto; padding: 38px 22px 80px; }}
    .topbar {{ display:flex; align-items:center; justify-content:space-between; gap: 16px; }}
    .chip {{
      display:inline-flex; align-items:center; gap:10px;
      padding: 10px 14px;
      border-radius: 999px;
      background: rgba(255,255,255,0.72);
      border: 1px solid rgba(0,0,0,0.08);
      backdrop-filter: blur(10px);
    }}
    .chip b {{ letter-spacing: 0.08em; font-size: 12px; }}
    .chip span {{ font-size: 12px; color: var(--muted); }}
    .actions {{ display:flex; gap: 10px; }}
    .btn {{
      appearance:none; border:1px solid rgba(0,0,0,0.10);
      background: rgba(255,255,255,0.72);
      border-radius: 12px;
      padding: 10px 12px;
      font-weight: 600;
      color: var(--text);
      cursor: pointer;
    }}
    .btn:hover {{ border-color: rgba(0,0,0,0.18); }}

    header {{ padding: 26px 0 18px; }}
    h1 {{ margin: 12px 0 8px; font-size: 44px; letter-spacing: -0.02em; line-height: 1.05; }}
    .subtitle {{ margin: 0; color: var(--muted); font-size: 17px; max-width: 70ch; }}

    .grid {{ display:grid; grid-template-columns: 1fr; gap: 18px; margin-top: 18px; }}
    @media (min-width: 980px) {{
      .grid {{ grid-template-columns: 1.15fr 0.85fr; }}
    }}

    .card {{
      background: rgba(255,255,255,0.78);
      border: 1px solid rgba(0,0,0,0.08);
      border-radius: var(--radius);
      box-shadow: var(--shadow);
      overflow: hidden;
    }}
    .card__hd {{ padding: 16px 18px; border-bottom: 1px solid rgba(0,0,0,0.06); }}
    .card__hd h2 {{ margin: 0; font-size: 16px; letter-spacing: 0.01em; }}
    .card__bd {{ padding: 16px 18px; }}
    .muted {{ color: var(--muted); }}

    .kpis {{ display:grid; grid-template-columns: repeat(2, minmax(0, 1fr)); gap: 12px; }}
    @media (min-width: 980px) {{
      .kpis {{ grid-template-columns: repeat(3, minmax(0, 1fr)); }}
    }}
    .kpi {{
      border-radius: 16px;
      background: rgba(255,255,255,0.78);
      border: 1px solid rgba(0,0,0,0.08);
      padding: 14px 14px 12px;
    }}
    .kpi__v {{ font-size: 26px; font-weight: 800; letter-spacing: -0.02em; }}
    .kpi__l {{ margin-top: 4px; font-size: 12px; color: var(--muted); }}
    .kpi--accent {{ border-color: color-mix(in oklab, var(--a1) 35%, rgba(0,0,0,0.08)); }}

    .section {{ margin-top: 18px; }}
    .section h3 {{ margin: 0 0 10px; font-size: 15px; letter-spacing: 0.01em; }}
    .bullets {{ margin: 0; padding-left: 18px; color: var(--muted); }}
    .bullets li {{ margin: 7px 0; }}

    .figure {{ padding: 12px 14px 14px; }}
    .figure img {{ width: 100%; height: auto; display:block; border-radius: 14px; border: 1px solid rgba(0,0,0,0.08); }}
    .cap {{ margin-top: 10px; font-size: 12px; color: var(--muted); }}

    .events {{ display:grid; gap: 10px; }}
    .event {{
      display:grid;
      grid-template-columns: 84px 1fr 44px;
      gap: 10px;
      padding: 12px;
      border-radius: 14px;
      background: rgba(255,255,255,0.72);
      border: 1px solid rgba(0,0,0,0.08);
      align-items: start;
    }}
    .event__when {{ font-weight: 700; font-size: 12px; color: var(--muted); }}
    .event__title {{ font-weight: 700; font-size: 13px; }}
    .event__badge {{
      justify-self: end;
      font-weight: 800;
      width: 32px; height: 32px;
      border-radius: 10px;
      display:flex; align-items:center; justify-content:center;
      background: color-mix(in oklab, var(--a1) 18%, white);
      border: 1px solid color-mix(in oklab, var(--a1) 25%, rgba(0,0,0,0.08));
      color: color-mix(in oklab, var(--text) 85%, black);
    }}
    .event__desc {{ grid-column: 2 / 4; font-size: 12px; color: var(--muted); }}

    .note {{
      padding: 14px 16px;
      border-radius: 16px;
      background: rgba(255,255,255,0.60);
      border: 1px solid rgba(0,0,0,0.08);
      color: var(--muted);
      font-size: 12px;
    }}

    details {{
      background: rgba(255,255,255,0.78);
      border: 1px solid rgba(0,0,0,0.08);
      border-radius: 16px;
      padding: 12px 14px;
    }}
    summary {{ cursor: pointer; font-weight: 700; }}
    .tablewrap {{ overflow:auto; border-radius: 14px; border: 1px solid rgba(0,0,0,0.08); margin-top: 12px; }}
    table {{ width: 100%; border-collapse: collapse; min-width: 920px; background: rgba(255,255,255,0.72); }}
    thead th {{ position: sticky; top: 0; background: rgba(255,255,255,0.92); border-bottom: 1px solid rgba(0,0,0,0.08); }}
    th, td {{ text-align: left; padding: 10px 10px; font-size: 12px; border-bottom: 1px solid rgba(0,0,0,0.06); vertical-align: top; }}
    tr:hover td {{ background: rgba(0,0,0,0.02); }}

    @media print {{
      body {{ background: white; }}
      .btn {{ display:none; }}
      .card {{ box-shadow:none; }}
      .chip {{ background: white; }}
    }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="topbar">
      <div class="chip"><b>INTEIA</b><span>Relatorio visual | Claude style</span></div>
      <div class="actions">
        {actions}
        <button class="btn" onclick="window.print()">Imprimir / PDF</button>
      </div>
    </div>

    <header>
      <h1>Celina Leao 2024-2025</h1>
      <p class="subtitle">Resultado consolidado das pesquisas (Jul/2024 a Dez/2025), com leitura de marcos e mensuracao do impacto do investimento em comunicacao a partir de Jan/2025.</p>
    </header>

    <div class="grid">
      <div class="card">
        <div class="card__hd"><h2>KPIs</h2></div>
        <div class="card__bd">
          <div class="kpis">
            <div class="kpi"><div class="kpi__v">{m.media_geral:.1f}%</div><div class="kpi__l">Media geral</div></div>
            <div class="kpi kpi--accent"><div class="kpi__v">+{m.lift_pp:.1f} p.p.</div><div class="kpi__l">Lift medio pos-MKT</div></div>
            <div class="kpi"><div class="kpi__v">{m.maximo:.0f}%</div><div class="kpi__l">Pico ({m.data_max})</div></div>
            <div class="kpi"><div class="kpi__v">{m.minimo:.0f}%</div><div class="kpi__l">Minimo ({m.data_min})</div></div>
            <div class="kpi"><div class="kpi__v">{m.total_pesquisas}</div><div class="kpi__l">Registros</div></div>
            <div class="kpi"><div class="kpi__v">{m.total_institutos}</div><div class="kpi__l">Institutos</div></div>
          </div>
          <div class="section" style="margin-top:14px;">
            <h3>Conclusao principal</h3>
            <ul class="bullets">
              <li><b>Hipotese confirmada no agregado:</b> apos Jan/2025, a media sobe de {m.media_pre:.1f}% para {m.media_pos:.1f}% (≈ {m.lift_rel*100:.0f}% de crescimento).</li>
              <li><b>Leitura:</b> sinal consistente de crescimento de marca, em sincronia com comunicacao mais intensa + agenda de entregas + eventos politicos/juridicos.</li>
              <li><b>Limite:</b> analise correlacional (institutos/cenarios variam). Para causalidade, recomenda-se tracking padronizado.</li>
            </ul>
          </div>
        </div>
      </div>

      <div class="card">
        <div class="card__hd"><h2>Marcos (linha do tempo)</h2></div>
        <div class="card__bd">
          <div class="events">{eventos_html}</div>
        </div>
      </div>
    </div>

    <div class="section">
      <div class="card">
        <div class="card__hd"><h2>Evolucao e evidencia</h2></div>
        <div class="card__bd">
          <div class="figure">
            <img src="{chart_imgs['evolucao']}" alt="Evolucao" />
            <div class="cap">Serie temporal (cenario 1). Sombreamento: periodo pos-MKT (>= Jan/2025).</div>
          </div>
          <div class="figure">
            <img src="{chart_imgs['pre_pos']}" alt="Pre pos" />
            <div class="cap">Comparacao pre vs pos com barras de dispersao (desvio-padrao). Lift: +{m.lift_pp:.1f} p.p.</div>
          </div>
          <div class="figure">
            <img src="{chart_imgs['tendencia']}" alt="Tendencia" />
            <div class="cap">Mudanca de inclinacao: regressao separada antes/depois sugere aceleracao apos Jan/2025.</div>
          </div>
          <div class="figure">
            <img src="{chart_imgs['waterfall']}" alt="Waterfall" />
            <div class="cap">Waterfall (cenario 1): evidencia de saltos (ganhos) e recuos (crises/concorrencia) ao longo do tempo.</div>
          </div>
        </div>
      </div>
    </div>

    <div class="section">
      <div class="card">
        <div class="card__hd"><h2>Triangulacao (institutos e cenario)</h2></div>
        <div class="card__bd">
          <div class="figure">
            <img src="{chart_imgs['institutos']}" alt="Institutos" />
            <div class="cap">Ultima medicao por instituto (datas diferentes). Use como triangulacao, nao como comparacao direta.</div>
          </div>
          <div class="figure">
            <img src="{chart_imgs['competicao']}" alt="Competicao" />
            <div class="cap">Snapshot do ultimo registro do dataset (nomes/cenarios variam conforme instituto).</div>
          </div>
          <div class="figure">
            <img src="{chart_imgs['boxplot']}" alt="Boxplot" />
            <div class="cap">Boxplot: distribuicao da intencao de voto por instituto (mede consistencia/variabilidade).</div>
          </div>
        </div>
      </div>
    </div>

    <div class="section">
      <div class="card">
        <div class="card__hd"><h2>Impacto da publicidade: interpretacao</h2></div>
        <div class="card__bd">
          <ul class="bullets">
            <li><b>Mensuracao:</b> lift medio pos-MKT (+{m.lift_pp:.1f} p.p.) e elevacao do teto observado (pico {m.maximo:.0f}%).</li>
            <li><b>Mecanismo provavel:</b> aumento de lembranca + repeticao de narrativa + conversao de visibilidade em preferencia.</li>
            <li><b>Marcos:</b> vitorias politicas (FCDF) e juridicas (absolvicao) agem como validadores; crises juridicas geram volatilidade.</li>
            <li><b>Risco de leitura:</b> entrada/saida de nomes (ex.: Fred/Arruda) desloca percentuais; por isso a serie usa cenario 1 como referencia.</li>
          </ul>
          <div class="figure">
            <img src="{chart_imgs['timeline']}" alt="Timeline" />
            <div class="cap">Eventos de alto impacto posicionados ao lado de pontos de pesquisa.</div>
          </div>
          <div class="note">
            Base: dados/pesquisas_eleitorais.csv e dados/eventos_timeline.json. Corte MKT: 01/01/2025.
          </div>
        </div>
      </div>
    </div>

    <div class="section">
      <div class="card">
        <div class="card__hd"><h2>Conclusao e recomendacoes</h2></div>
        <div class="card__bd">
          <ul class="bullets">
            <li><b>Conclusao:</b> o padrao agregado e consistente com crescimento de marca apos Jan/2025, com consolidacao e pico em Set/2025.</li>
            <li><b>Recomendacao 1:</b> manter cadencia de comunicacao (press + digital) e mensagens repetiveis.</li>
            <li><b>Recomendacao 2:</b> criar tracking padronizado e metricas de marca (share of voice, sentimento, busca) para medir impacto com precisao.</li>
            <li><b>Recomendacao 3:</b> gestao ativa de risco juridico e resposta rapida a ataques, para evitar erosao em janelas criticas.</li>
          </ul>
        </div>
      </div>
    </div>

    <div class="section">
      <details>
        <summary>Base de pesquisas (CSV) - ver tabela</summary>
        <div class="muted" style="margin-top:8px; font-size:12px;">Fonte: <code>dados/pesquisas_eleitorais.csv</code>. Cada linha e um cenario divulgado por instituto.</div>
        <div class="tablewrap">
          <table>
            <thead>
              <tr>
                <th>Divulgacao</th>
                <th>Instituto</th>
                <th>Cenario</th>
                <th>Celina</th>
                <th>2o</th>
                <th>3o</th>
                <th>Obs.</th>
              </tr>
            </thead>
            <tbody>
              {tabela_html}
            </tbody>
          </table>
        </div>
      </details>
    </div>

  </div>
</body>
</html>
"""

    out = out_path or (OUT_DIR / "index.html")
    out.write_text(html, encoding="utf-8")
    return out


def main() -> int:
    tokens = carregar_tokens(ROOT / "presentation_sdk" / "tokens" / "anthropic_light.json")
    df, eventos = carregar_dados()
    m = calcular_metricas(df)

    CHARTS_DIR.mkdir(parents=True, exist_ok=True)
    charts = gerar_graficos(df, eventos, m, tokens)

    pptx = gerar_pptx(df, eventos, m, charts, tokens)
    md = gerar_marp(m, charts)
    html = gerar_html(df, eventos, m, charts, tokens, download_links=True, out_path=OUT_DIR / "index.html")

    relatorio = ROOT / "relatorio" / "relatorio-claude.html"
    relatorio.parent.mkdir(parents=True, exist_ok=True)
    gerar_html(df, eventos, m, charts, tokens, download_links=False, out_path=relatorio)

    print(str(pptx))
    print(str(md))
    print(str(html))
    print(str(relatorio))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
