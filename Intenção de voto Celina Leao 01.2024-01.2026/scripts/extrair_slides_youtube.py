#!/usr/bin/env python3
"""Extrai slides (prints) de videos do YouTube.

Pipeline:
1) Baixa o video via yt-dlp (mp4 quando possivel)
2) Extrai frames em mudancas de cena (ffmpeg scene detection)
3) Deduplica frames parecidos (perceptual hash)
4) Gera:
   - imagens finais: slide_###.png
   - contact_sheet.jpg (visao geral)
   - index.html (galeria simples, com zoom do navegador)
   - deck.pptx (um slide por imagem, 16:9)
   - style.json (paleta de cores estimada)

Requisitos:
- ffmpeg no PATH
- Python libs: yt-dlp, pillow, imagehash, colorthief, python-pptx
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from colorthief import ColorThief
from PIL import Image, ImageOps
import imagehash


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUT_DIR = ROOT / "outputs" / "youtube_slides"


@dataclass(frozen=True)
class VideoInfo:
    video_id: str
    title: str
    webpage_url: str


def _run(cmd: list[str]) -> None:
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"Falhou: {' '.join(cmd)}\n\n{proc.stdout}")


def _run_capture(cmd: list[str]) -> str:
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"Falhou: {' '.join(cmd)}\n\n{proc.stdout}")
    return proc.stdout


def _ensure_tool(tool: str) -> None:
    if _tool_path(tool) is None:
        raise RuntimeError(f"Ferramenta ausente no PATH/venv: {tool}")


def _tool_path(tool: str) -> str | None:
    p = shutil.which(tool)
    if p:
        return p
    # Quando roda via venv, os entrypoints ficam em <sys.prefix>/bin.
    venv_candidate = Path(sys.prefix).resolve() / "bin" / tool
    if venv_candidate.exists():
        return str(venv_candidate)
    # Fallback: mesma pasta do executavel.
    exe_candidate = Path(sys.executable).resolve().parent / tool
    if exe_candidate.exists():
        return str(exe_candidate)
    return None


def _slugify(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text or "video"


def get_video_info(url: str) -> VideoInfo:
    _ensure_tool("yt-dlp")
    yt_dlp = _tool_path("yt-dlp")
    assert yt_dlp is not None
    out = _run_capture(
        [
            yt_dlp,
            "--no-warnings",
            "--skip-download",
            "--print",
            "%(.id)s\t%(.title)s\t%(.webpage_url)s",
            url,
        ]
    ).strip()
    if not out:
        raise RuntimeError("Nao consegui ler metadados do video.")
    parts = out.split("\t")
    if len(parts) < 3:
        raise RuntimeError(f"Formato inesperado de metadados: {out}")
    return VideoInfo(video_id=parts[0], title=parts[1], webpage_url=parts[2])


def _looks_like_single_video(url: str) -> bool:
    if "watch?v=" in url:
        return True
    if "youtu.be/" in url:
        return True
    if "/shorts/" in url:
        return True
    return False


def list_videos(url: str, *, limit: int) -> list[str]:
    _ensure_tool("yt-dlp")
    yt_dlp = _tool_path("yt-dlp")
    assert yt_dlp is not None

    # Para canais, o tab /videos respeita melhor playlist-end.
    url = url.strip()
    if "youtube.com/@" in url or "/channel/" in url:
        url = url.rstrip("/")
        if not url.endswith("/videos"):
            url = url + "/videos"

    out = _run_capture(
        [
            yt_dlp,
            "--no-warnings",
            "--flat-playlist",
            "--playlist-end",
            str(limit),
            "--print",
            "%(webpage_url)s",
            url,
        ]
    )
    urls = [ln.strip() for ln in out.splitlines() if ln.strip()]
    # Filtra shorts (normalmente nao sao apresentacoes de slides)
    urls = [u for u in urls if "/shorts/" not in u]
    return urls


def download_video(url: str, dest_mp4: Path) -> None:
    _ensure_tool("yt-dlp")
    yt_dlp = _tool_path("yt-dlp")
    assert yt_dlp is not None
    dest_mp4.parent.mkdir(parents=True, exist_ok=True)
    # O YouTube as vezes bloqueia alguns formatos (403). Esta funcao tenta
    # uma ordem de fallbacks para maximizar chance e manter qualidade.
    candidates: list[list[str]] = [
        # Prefere AVC (muito menos propenso a 403 do que AV1 em alguns videos)
        [
            "-f",
            "bv*[ext=mp4][vcodec^=avc1]+ba[ext=m4a]/b[ext=mp4][vcodec^=avc1]/b[ext=mp4]/b",
        ],
        # Fallback generico
        ["-f", "bv*[ext=mp4]+ba[ext=m4a]/b[ext=mp4]/b"],
        # Ultimo recurso: MP4 progressivo (tipicamente 360p) via player android
        ["--extractor-args", "youtube:player_client=android", "-f", "18"],
    ]

    last_out = ""
    for extra in candidates:
        cmd = [
            yt_dlp,
            "--no-warnings",
            "--no-playlist",
            "-o",
            str(dest_mp4),
            *extra,
            url,
        ]
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
        if proc.returncode == 0:
            return
        last_out = proc.stdout

    raise RuntimeError(f"Falhou: download do video\n\n{last_out}")


def extract_scene_frames(
    input_mp4: Path,
    frames_dir: Path,
    *,
    scene_threshold: float,
    max_width: int,
    sharpen: bool,
) -> None:
    _ensure_tool("ffmpeg")
    frames_dir.mkdir(parents=True, exist_ok=True)

    # Seleciona o primeiro frame e os frames onde o detector de cena acusa mudanca.
    # A escala serve para "dar zoom" (aumentar resolucao) e melhorar legibilidade.
    # OBS: max_width preserva aspect ratio.
    vf_parts = [
        f"select='eq(n,0)+gt(scene,{scene_threshold})'",
        f"scale='min(iw,{max_width})':-2:flags=lanczos",
    ]
    if sharpen:
        vf_parts.append("unsharp=5:5:1.0:3:3:0.0")
    vf = ",".join(vf_parts)

    out_pattern = str(frames_dir / "raw_%06d.png")
    _run(
        [
            "ffmpeg",
            "-hide_banner",
            "-loglevel",
            "error",
            "-i",
            str(input_mp4),
            "-vf",
            vf,
            "-vsync",
            "vfr",
            "-q:v",
            "2",
            out_pattern,
        ]
    )


def dedupe_frames(
    frames_dir: Path,
    slides_dir: Path,
    *,
    phash_distance: int,
    min_keep: int,
    crop_screen: bool,
    crop_min_area: float,
    crop_target_width: int,
) -> list[Path]:
    slides_dir.mkdir(parents=True, exist_ok=True)
    raw = sorted(frames_dir.glob("raw_*.png"))
    if not raw:
        raise RuntimeError("Nenhum frame extraido. Ajuste o scene_threshold.")

    kept: list[Path] = []
    last_hash = None
    slide_idx = 1

    for p in raw:
        with Image.open(p) as img0:
            img0 = ImageOps.exif_transpose(img0)
            img = img0
            if crop_screen:
                cropped = crop_screen_image(img0, min_area_ratio=crop_min_area, target_width=crop_target_width)
                if cropped is None:
                    continue
                img = cropped
            h = imagehash.phash(img)
        if last_hash is None or (h - last_hash) > phash_distance:
            out = slides_dir / f"slide_{slide_idx:03d}.png"
            if crop_screen:
                img.save(out)
            else:
                shutil.copyfile(p, out)
            kept.append(out)
            last_hash = h
            slide_idx += 1

    # Garante um minimo (quando a apresentacao tem muitas animacoes e hashes ficam parecidos).
    if len(kept) < min_keep:
        # Faz uma segunda passada, mais permissiva.
        for p in raw:
            if len(kept) >= min_keep:
                break
            with Image.open(p) as img0:
                img0 = ImageOps.exif_transpose(img0)
                img = img0
                if crop_screen:
                    cropped = crop_screen_image(
                        img0, min_area_ratio=crop_min_area, target_width=crop_target_width
                    )
                    if cropped is None:
                        continue
                    img = cropped
                h = imagehash.phash(img)
            if last_hash is None or (h - last_hash) > max(1, phash_distance // 2):
                out = slides_dir / f"slide_{slide_idx:03d}.png"
                if not out.exists():
                    if crop_screen:
                        img.save(out)
                    else:
                        shutil.copyfile(p, out)
                    kept.append(out)
                    last_hash = h
                    slide_idx += 1

    return kept


def crop_screen_image(
    img: Image.Image,
    *,
    min_area_ratio: float,
    target_width: int,
) -> Image.Image | None:
    """Tenta detectar a tela do palco e recortar/retificar.

    Funciona melhor para keynotes gravadas (camera + telona). Para videos que ja
    sao apenas os slides, normalmente retorna a imagem inteira.
    """

    try:
        import cv2  # type: ignore
        import numpy as np
    except Exception:
        return None

    rgb = np.array(img.convert("RGB"))
    bgr = rgb[:, :, ::-1]

    h_img, w_img = bgr.shape[:2]
    img_area = float(h_img * w_img)

    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
    gray = cv2.GaussianBlur(gray, (7, 7), 0)
    edges = cv2.Canny(gray, 40, 140)
    edges = cv2.dilate(edges, None, iterations=1)
    edges = cv2.erode(edges, None, iterations=1)

    cnts, _hier = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    if not cnts:
        return None

    cnts = sorted(cnts, key=cv2.contourArea, reverse=True)[:12]
    best = None
    best_area = 0.0

    for c in cnts:
        area = float(cv2.contourArea(c))
        if area < img_area * min_area_ratio:
            continue
        peri = cv2.arcLength(c, True)
        approx = cv2.approxPolyDP(c, 0.02 * peri, True)
        if len(approx) == 4:
            if area > best_area:
                best = approx
                best_area = area

    if best is None:
        # Fallback: bounding box do maior contorno (sem perspectiva)
        c = cnts[0]
        area = float(cv2.contourArea(c))
        if area < img_area * min_area_ratio:
            return None
        x, y, w, h = cv2.boundingRect(c)
        crop = bgr[y : y + h, x : x + w]
        out = Image.fromarray(crop[:, :, ::-1])
        return _resize_up_to(out, target_width)

    pts = best.reshape(4, 2).astype("float32")
    warp = _four_point_transform(bgr, pts)
    if warp is None:
        return None

    out = Image.fromarray(warp[:, :, ::-1])
    return _resize_up_to(out, target_width)


def _resize_up_to(img: Image.Image, target_width: int) -> Image.Image:
    if target_width <= 0:
        return img
    w, h = img.size
    if w >= target_width:
        return img
    scale = target_width / float(w)
    return img.resize((int(w * scale), int(h * scale)), Image.Resampling.LANCZOS)


def _order_points(pts):
    import numpy as np

    rect = np.zeros((4, 2), dtype="float32")
    s = pts.sum(axis=1)
    rect[0] = pts[np.argmin(s)]
    rect[2] = pts[np.argmax(s)]

    diff = np.diff(pts, axis=1)
    rect[1] = pts[np.argmin(diff)]
    rect[3] = pts[np.argmax(diff)]
    return rect


def _four_point_transform(image_bgr, pts):
    import cv2  # type: ignore
    import numpy as np

    rect = _order_points(pts)
    (tl, tr, br, bl) = rect

    width_a = np.linalg.norm(br - bl)
    width_b = np.linalg.norm(tr - tl)
    max_w = int(max(width_a, width_b))

    height_a = np.linalg.norm(tr - br)
    height_b = np.linalg.norm(tl - bl)
    max_h = int(max(height_a, height_b))

    if max_w < 20 or max_h < 20:
        return None

    dst = np.array(
        [[0, 0], [max_w - 1, 0], [max_w - 1, max_h - 1], [0, max_h - 1]], dtype="float32"
    )
    m = cv2.getPerspectiveTransform(rect, dst)
    return cv2.warpPerspective(image_bgr, m, (max_w, max_h))


def _rgb_to_hex(rgb: tuple[int, int, int]) -> str:
    return "#%02x%02x%02x" % rgb


def estimate_palette(images: Iterable[Path], *, quality: int = 10) -> dict:
    # Agrega cores dominantes por slide. Nao e perfeito, mas gera uma base boa.
    counts: dict[tuple[int, int, int], int] = {}
    for p in images:
        try:
            ct = ColorThief(str(p))
            pal = ct.get_palette(color_count=6, quality=quality)
        except Exception:
            continue
        for c in pal:
            counts[c] = counts.get(c, 0) + 1

    ordered = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)
    top = [rgb for rgb, _ in ordered[:12]]
    hexes = [_rgb_to_hex(c) for c in top]

    # Heuristica simples: supoe que o fundo e a cor mais frequente.
    background = hexes[0] if hexes else "#ffffff"
    accent = hexes[1] if len(hexes) > 1 else "#111111"

    return {
        "cores": {
            "background": background,
            "accent": accent,
            "palette": hexes,
        },
        "tipografia": {
            "fonte_base": "(defina aqui)",
            "fallback": ["Inter", "Helvetica Neue", "Arial"],
        },
        "graficos": {
            "grid": "sutil",
            "eixos": "minimos",
            "legendas": "curtas",
        },
    }


def make_contact_sheet(images: list[Path], out_path: Path, *, cols: int = 5, thumb_w: int = 420) -> None:
    if not images:
        return
    thumbs: list[Image.Image] = []
    for p in images:
        img = Image.open(p)
        img = ImageOps.exif_transpose(img)
        w, h = img.size
        scale = thumb_w / float(w)
        img = img.resize((thumb_w, int(h * scale)), Image.Resampling.LANCZOS)
        thumbs.append(img)

    rows = int(math.ceil(len(thumbs) / cols))
    pad = 12
    cell_w = thumb_w
    cell_h = max(t.size[1] for t in thumbs)
    sheet = Image.new("RGB", (cols * cell_w + (cols + 1) * pad, rows * cell_h + (rows + 1) * pad), (245, 245, 245))

    for i, t in enumerate(thumbs):
        r = i // cols
        c = i % cols
        x = pad + c * (cell_w + pad)
        y = pad + r * (cell_h + pad)
        sheet.paste(t, (x, y))
        t.close()

    out_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path, quality=92)


def make_gallery(images: list[Path], out_path: Path, *, title: str) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    rels = [p.name for p in images]
    html = f"""<!doctype html>
<html lang=\"pt-br\">
<meta charset=\"utf-8\" />
<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
<title>{title}</title>
<style>
  :root {{ --bg:#0b0d10; --card:#11151c; --text:#e9eef7; --muted:#9aa7b7; }}
  body {{ margin:0; font-family: Inter, system-ui, -apple-system, Segoe UI, Roboto, Arial, sans-serif; background:linear-gradient(180deg, #0b0d10, #06070a); color:var(--text); }}
  header {{ padding:24px 20px 10px; }}
  h1 {{ margin:0; font-size:18px; letter-spacing:0.2px; }}
  p {{ margin:8px 0 0; color:var(--muted); font-size:13px; }}
  .grid {{ display:grid; grid-template-columns: repeat(auto-fill, minmax(260px, 1fr)); gap:14px; padding:14px 20px 28px; }}
  a {{ display:block; background:var(--card); border:1px solid rgba(255,255,255,0.06); border-radius:14px; overflow:hidden; text-decoration:none; }}
  img {{ width:100%; height:auto; display:block; }}
  .cap {{ padding:10px 12px; font-size:12px; color:var(--muted); }}
</style>
<header>
  <h1>{title}</h1>
  <p>Clique em um slide para abrir em tamanho total (zoom pelo navegador).</p>
</header>
<div class=\"grid\">
"""
    for i, name in enumerate(rels, start=1):
        html += f"  <a href=\"{name}\" target=\"_blank\" rel=\"noreferrer\"><img src=\"{name}\" alt=\"Slide {i}\" loading=\"lazy\" /><div class=\"cap\">Slide {i:03d}</div></a>\n"
    html += """</div>
</html>
"""
    out_path.write_text(html, encoding="utf-8")


def make_pptx(images: list[Path], out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    for p in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Extrai prints de slides de um video do YouTube")
    parser.add_argument("--url", required=True, help="URL do video do YouTube")
    parser.add_argument("--limit", type=int, default=5, help="Quando --url for canal/playlist: quantos videos processar")
    parser.add_argument("--out", default=str(DEFAULT_OUT_DIR), help="Pasta base de saida")
    parser.add_argument("--scene-threshold", type=float, default=0.22, help="Sensibilidade do detector de cena (menor = mais frames)")
    parser.add_argument("--max-width", type=int, default=2560, help="Largura maxima do print (faz upscale/zoom ate esse limite)")
    parser.add_argument("--sharpen", action="store_true", help="Aplica um leve sharpening nos prints")
    parser.add_argument("--phash-distance", type=int, default=6, help="Distancia minima de hash para considerar frame novo")
    parser.add_argument("--min-keep", type=int, default=10, help="Mantem pelo menos N slides")
    parser.add_argument(
        "--crop-screen",
        action="store_true",
        help="Tenta detectar/recortar a tela (keynotes gravadas com palco)",
    )
    parser.add_argument(
        "--crop-min-area",
        type=float,
        default=0.18,
        help="Area minima (0-1) para aceitar a tela detectada",
    )
    args = parser.parse_args(argv)

    urls = [args.url]
    if not _looks_like_single_video(args.url):
        urls = list_videos(args.url, limit=args.limit)
        if not urls:
            raise RuntimeError("Nao encontrei videos na URL informada.")

    bases: list[str] = []
    for u in urls:
        info = get_video_info(u)
        safe = _slugify(info.title)
        base = Path(args.out) / f"{info.video_id}_{safe}"
        raw_dir = base / "raw"
        frames_dir = base / "frames"
        suffix = "_screen" if args.crop_screen else ""
        slides_dir = base / ("slides_screen" if args.crop_screen else "slides")

        mp4 = raw_dir / f"{info.video_id}.mp4"
        if not mp4.exists():
            download_video(info.webpage_url, mp4)

        extract_scene_frames(
            mp4,
            frames_dir,
            scene_threshold=args.scene_threshold,
            max_width=args.max_width,
            sharpen=args.sharpen,
        )

        slides = dedupe_frames(
            frames_dir,
            slides_dir,
            phash_distance=args.phash_distance,
            min_keep=args.min_keep,
            crop_screen=args.crop_screen,
            crop_min_area=args.crop_min_area,
            crop_target_width=args.max_width,
        )

        make_contact_sheet(slides, base / f"contact_sheet{suffix}.jpg")
        make_gallery(slides, base / f"index{suffix}.html", title=info.title)
        make_pptx(slides, base / f"deck{suffix}.pptx")

        style = estimate_palette(slides)
        style.update(
            {
                "video": {
                    "id": info.video_id,
                    "titulo": info.title,
                    "url": info.webpage_url,
                },
                "params": {
                    "scene_threshold": args.scene_threshold,
                    "max_width": args.max_width,
                    "phash_distance": args.phash_distance,
                },
            }
        )
        (base / f"style{suffix}.json").write_text(
            json.dumps(style, indent=2, ensure_ascii=True), encoding="utf-8"
        )
        bases.append(str(base))

    if len(bases) > 1:
        batch_path = Path(args.out) / "batch_index.json"
        batch_path.write_text(json.dumps({"items": bases}, indent=2, ensure_ascii=True), encoding="utf-8")

    print("\n".join(bases))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
